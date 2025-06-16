from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime

class PitchDeckGenerator:
    def __init__(self):
        self.prs = Presentation()
        self._setup_slide_layouts()
        
    def _setup_slide_layouts(self):
        """Define custom slide layouts and styles"""
        # Define colors
        self.PRIMARY_COLOR = RGBColor(0, 112, 192)  # Professional blue
        self.SECONDARY_COLOR = RGBColor(68, 68, 68)  # Dark gray
        self.ACCENT_COLOR = RGBColor(255, 102, 0)   # Orange accent
        
        # Define fonts
        self.TITLE_FONT = 'Calibri'
        self.BODY_FONT = 'Calibri'
        
    def _create_title_slide(self, startup_name, tagline):
        """Create the cover slide with professional styling"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        
        # Add title
        title = slide.shapes.title
        title.text = startup_name
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add tagline
        subtitle = slide.placeholders[1]
        subtitle.text = tagline
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.SECONDARY_COLOR
        
        # Add date
        date_box = slide.shapes.add_textbox(Inches(8), Inches(6.5), Inches(2), Inches(0.5))
        date_box.text = datetime.now().strftime("%B %Y")
        date_box.text_frame.paragraphs[0].font.size = Pt(12)
        date_box.text_frame.paragraphs[0].font.color.rgb = self.SECONDARY_COLOR
        
    def _create_content_slide(self, title, content, section_type):
        """Create a content slide with consistent styling"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        
        # Process content based on section type
        if section_type in ['problem', 'solution', 'market']:
            # Use bullet points
            for point in content.split('\n'):
                if point.strip():
                    p = tf.add_paragraph()
                    p.text = point.strip()
                    p.font.size = Pt(18)
                    p.font.color.rgb = self.SECONDARY_COLOR
                    p.level = 0
        else:
            # Use regular text
            p = tf.add_paragraph()
            p.text = content
            p.font.size = Pt(18)
            p.font.color.rgb = self.SECONDARY_COLOR
            
    def generate_pitch_deck(self, form_data, generated_deck):
        """Generate a complete pitch deck from the provided data"""
        # Create cover slide using form_data
        self._create_title_slide(
            form_data['startup_name'],
            form_data.get('tagline', 'Transforming Ideas into Reality')
        )
        
        # Create content slides using generated_deck and SLIDE_SECTIONS order
        sections_map = {
            'cover': 'Cover Slide',
            'problem': 'The Problem',
            'solution': 'Our Solution',
            'market': 'Market Opportunity',
            'product': 'Product Overview',
            'business_model': 'Business Model',
            'competition': 'Competitive Advantage',
            'team': 'Our Team',
            'traction': 'Traction & Milestones',
            'funding_needs': 'Investment Opportunity'
        }
        
        for section_key in sections_map.keys():
            title = sections_map[section_key]
            content = generated_deck.get(section_key)
            
            if content:
                self._create_content_slide(
                    title,
                    content,
                    section_key
                )
                
    def save(self, filename):
        """Save the presentation to a file"""
        self.prs.save(filename)
        return os.path.abspath(filename)

def generate_ppt(form_data, generated_deck, output_dir='generated_decks'):
    """Generate a PowerPoint presentation from the pitch deck data"""
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Generate filename using startup_name from form_data
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{form_data['startup_name'].replace(' ', '_')}_{timestamp}.pptx"
    filepath = os.path.join(output_dir, filename)
    
    # Generate the presentation
    generator = PitchDeckGenerator()
    generator.generate_pitch_deck(form_data, generated_deck)
    
    # Save and return the file path
    return generator.save(filepath) 